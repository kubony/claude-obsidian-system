#!/usr/bin/env python3
"""
PPT Generator - 한국어 최적화 미니멀 프레젠테이션 생성

python-pptx를 사용하여 PPTX 파일을 직접 생성합니다.
"""

import argparse
import json
import os
import sys
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# 컬러 팔레트 정의
PALETTES = {
    "sage": {
        "surface": "#f5f5f0",
        "surface_foreground": "#1a1a1a",
        "primary": "#b8c4b8",
        "accent": "#2d2d2d",
        "muted": "#e8e8e3",
        "muted_foreground": "#666666",
        "dark_bg": "#2d2d2d",
        "dark_fg": "#f5f5f0",
    },
    "mono": {
        "surface": "#ffffff",
        "surface_foreground": "#111111",
        "primary": "#f0f0f0",
        "accent": "#111111",
        "muted": "#f5f5f5",
        "muted_foreground": "#666666",
        "dark_bg": "#111111",
        "dark_fg": "#ffffff",
    },
    "navy": {
        "surface": "#f8f9fc",
        "surface_foreground": "#1a1f36",
        "primary": "#dce3f0",
        "accent": "#1a1f36",
        "muted": "#eef1f6",
        "muted_foreground": "#5c6478",
        "dark_bg": "#1a1f36",
        "dark_fg": "#f8f9fc",
    },
}

# 폰트 설정
FONTS = {
    "display": "Noto Serif KR",  # 제목용
    "content": "Pretendard",      # 본문용
    "fallback_display": "맑은 고딕",
    "fallback_content": "맑은 고딕",
}

# 슬라이드 크기 (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# 여백
MARGIN = Inches(0.75)


def hex_to_rgb(hex_color: str) -> RGBColor:
    """HEX 컬러를 RGBColor로 변환"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def set_text_frame_properties(text_frame, font_name: str, font_size: int,
                               font_color: str, bold: bool = False,
                               alignment: PP_ALIGN = PP_ALIGN.LEFT):
    """텍스트 프레임 속성 설정"""
    text_frame.word_wrap = True
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = alignment
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = hex_to_rgb(font_color)
        paragraph.font.bold = bold


def add_text_box(slide, left: float, top: float, width: float, height: float,
                 text: str, font_name: str, font_size: int, font_color: str,
                 bold: bool = False, alignment: PP_ALIGN = PP_ALIGN.LEFT,
                 vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP) -> None:
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = text
    tf.paragraphs[0].alignment = alignment
    tf.paragraphs[0].font.name = font_name
    tf.paragraphs[0].font.size = Pt(font_size)
    tf.paragraphs[0].font.color.rgb = hex_to_rgb(font_color)
    tf.paragraphs[0].font.bold = bold

    # 수직 정렬
    txBox.text_frame.auto_size = None
    txBox.text_frame.word_wrap = True


def add_shape_with_fill(slide, left: float, top: float, width: float,
                        height: float, fill_color: str) -> None:
    """배경 사각형 추가"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    shape.line.fill.background()


def create_cover_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """표지 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # 배경색
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 제목
    add_text_box(
        slide, 0.75, 2.5, 11.833, 1.5,
        data.get("title", "제목"),
        FONTS["display"], 44, palette["surface_foreground"],
        bold=True, alignment=PP_ALIGN.CENTER
    )

    # 부제목
    if data.get("subtitle"):
        add_text_box(
            slide, 0.75, 4.0, 11.833, 0.8,
            data["subtitle"],
            FONTS["content"], 20, palette["muted_foreground"],
            alignment=PP_ALIGN.CENTER
        )

    # 발표자 / 날짜
    footer_text = ""
    if data.get("author"):
        footer_text = data["author"]
    if data.get("date"):
        footer_text += f"  |  {data['date']}" if footer_text else data["date"]

    if footer_text:
        add_text_box(
            slide, 0.75, 6.2, 11.833, 0.5,
            footer_text,
            FONTS["content"], 14, palette["muted_foreground"],
            alignment=PP_ALIGN.CENTER
        )


def create_section_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """섹션 구분 슬라이드 생성 (다크 배경)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 다크 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["dark_bg"])
    background.line.fill.background()

    # 배지
    if data.get("badge"):
        add_text_box(
            slide, 0.75, 2.8, 11.833, 0.5,
            data["badge"],
            FONTS["content"], 14, palette["muted_foreground"],
            alignment=PP_ALIGN.CENTER
        )

    # 섹션 제목
    add_text_box(
        slide, 0.75, 3.3, 11.833, 1.2,
        data.get("title", "섹션 제목"),
        FONTS["display"], 36, palette["dark_fg"],
        bold=True, alignment=PP_ALIGN.CENTER
    )


def create_stats_grid_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """통계 그리드 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 슬라이드 제목
    add_text_box(
        slide, 0.75, 0.75, 11.833, 0.8,
        data.get("title", "핵심 지표"),
        FONTS["display"], 32, palette["surface_foreground"],
        bold=True
    )

    # 통계 항목들
    stats = data.get("stats", [])
    num_stats = len(stats)

    if num_stats > 0:
        # 그리드 계산
        col_width = 11.833 / min(num_stats, 4)
        start_y = 2.5

        for i, stat in enumerate(stats[:4]):  # 최대 4개
            col_x = 0.75 + (i * col_width)

            # 숫자
            add_text_box(
                slide, col_x, start_y, col_width - 0.3, 1.2,
                stat.get("number", "0"),
                FONTS["display"], 48, palette["accent"],
                bold=True, alignment=PP_ALIGN.CENTER
            )

            # 레이블
            add_text_box(
                slide, col_x, start_y + 1.3, col_width - 0.3, 0.8,
                stat.get("label", ""),
                FONTS["content"], 14, palette["muted_foreground"],
                alignment=PP_ALIGN.CENTER
            )


def create_two_column_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """2단 비교 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 슬라이드 제목
    add_text_box(
        slide, 0.75, 0.75, 11.833, 0.8,
        data.get("title", "비교 분석"),
        FONTS["display"], 32, palette["surface_foreground"],
        bold=True
    )

    col_width = 5.5

    # 왼쪽 컬럼
    left_data = data.get("left", {})
    add_text_box(
        slide, 0.75, 2.0, col_width, 0.6,
        left_data.get("heading", "왼쪽"),
        FONTS["display"], 24, palette["surface_foreground"],
        bold=True
    )

    left_items = left_data.get("items", [])
    for j, item in enumerate(left_items[:5]):
        add_text_box(
            slide, 0.75, 2.8 + (j * 0.6), col_width, 0.5,
            f"• {item}",
            FONTS["content"], 16, palette["surface_foreground"]
        )

    # 오른쪽 컬럼
    right_data = data.get("right", {})
    add_text_box(
        slide, 7.0, 2.0, col_width, 0.6,
        right_data.get("heading", "오른쪽"),
        FONTS["display"], 24, palette["surface_foreground"],
        bold=True
    )

    right_items = right_data.get("items", [])
    for j, item in enumerate(right_items[:5]):
        add_text_box(
            slide, 7.0, 2.8 + (j * 0.6), col_width, 0.5,
            f"• {item}",
            FONTS["content"], 16, palette["surface_foreground"]
        )


def create_three_column_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """3단 컬럼 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 슬라이드 제목
    add_text_box(
        slide, 0.75, 0.75, 11.833, 0.8,
        data.get("title", "핵심 기능"),
        FONTS["display"], 32, palette["surface_foreground"],
        bold=True
    )

    columns = data.get("columns", [])
    col_width = 3.8

    for i, col in enumerate(columns[:3]):
        col_x = 0.75 + (i * 4.1)

        # 컬럼 제목
        add_text_box(
            slide, col_x, 2.2, col_width, 0.6,
            col.get("heading", f"기능 {i+1}"),
            FONTS["display"], 20, palette["surface_foreground"],
            bold=True
        )

        # 컬럼 설명
        add_text_box(
            slide, col_x, 3.0, col_width, 2.5,
            col.get("description", ""),
            FONTS["content"], 14, palette["muted_foreground"]
        )


def create_image_text_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """이미지 + 텍스트 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 슬라이드 제목
    add_text_box(
        slide, 0.75, 0.75, 11.833, 0.8,
        data.get("title", "스토리텔링"),
        FONTS["display"], 32, palette["surface_foreground"],
        bold=True
    )

    # 이미지 영역 (왼쪽)
    image_path = data.get("image_path")
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(
                image_path,
                Inches(0.75), Inches(2.0),
                width=Inches(5.5)
            )
        except Exception as e:
            # 이미지 로드 실패 시 placeholder
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.75), Inches(2.0), Inches(5.5), Inches(4.0)
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = hex_to_rgb(palette["muted"])
            placeholder.line.fill.background()
    else:
        # 이미지 없을 때 placeholder
        placeholder = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.75), Inches(2.0), Inches(5.5), Inches(4.0)
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = hex_to_rgb(palette["muted"])
        placeholder.line.fill.background()

    # 텍스트 영역 (오른쪽)
    add_text_box(
        slide, 6.75, 2.0, 5.5, 4.0,
        data.get("text", "설명 텍스트"),
        FONTS["content"], 16, palette["surface_foreground"]
    )


def create_closing_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """마무리 슬라이드 생성"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 감사 메시지
    add_text_box(
        slide, 0.75, 2.8, 11.833, 1.2,
        data.get("title", "감사합니다"),
        FONTS["display"], 44, palette["surface_foreground"],
        bold=True, alignment=PP_ALIGN.CENTER
    )

    # 연락처
    if data.get("contact"):
        add_text_box(
            slide, 0.75, 4.5, 11.833, 0.6,
            data["contact"],
            FONTS["content"], 16, palette["muted_foreground"],
            alignment=PP_ALIGN.CENTER
        )


def create_content_slide(prs: Presentation, data: Dict, palette: Dict) -> None:
    """일반 콘텐츠 슬라이드 생성 (기본 레이아웃)"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 배경
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(palette["surface"])
    background.line.fill.background()

    # 슬라이드 제목
    add_text_box(
        slide, 0.75, 0.75, 11.833, 0.8,
        data.get("title", "제목"),
        FONTS["display"], 32, palette["surface_foreground"],
        bold=True
    )

    # 본문 내용
    content = data.get("content", "")
    if isinstance(content, list):
        content = "\n".join([f"• {item}" for item in content])

    add_text_box(
        slide, 0.75, 2.0, 11.833, 4.5,
        content,
        FONTS["content"], 16, palette["surface_foreground"]
    )


# 레이아웃 타입별 생성 함수 매핑
LAYOUT_CREATORS = {
    "cover": create_cover_slide,
    "section": create_section_slide,
    "stats_grid": create_stats_grid_slide,
    "two_column": create_two_column_slide,
    "three_column": create_three_column_slide,
    "image_text": create_image_text_slide,
    "closing": create_closing_slide,
    "content": create_content_slide,
}


def generate_pptx(config: Dict, output_path: str, palette_name: str = "sage") -> str:
    """PPTX 파일 생성 메인 함수"""

    # 팔레트 선택
    palette = PALETTES.get(palette_name, PALETTES["sage"])

    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 슬라이드 생성
    slides_data = config.get("slides", [])

    for slide_data in slides_data:
        layout = slide_data.get("layout", "content")
        creator_func = LAYOUT_CREATORS.get(layout, create_content_slide)
        creator_func(prs, slide_data, palette)

    # 파일 저장
    output_path = os.path.expanduser(output_path)
    prs.save(output_path)

    return output_path


def main():
    parser = argparse.ArgumentParser(
        description="한국어 최적화 미니멀 PPTX 생성기"
    )
    parser.add_argument(
        "--config", "-c",
        required=True,
        help="슬라이드 설정 JSON 파일 경로"
    )
    parser.add_argument(
        "--output", "-o",
        default="output.pptx",
        help="출력 PPTX 파일 경로 (기본: output.pptx)"
    )
    parser.add_argument(
        "--palette", "-p",
        choices=["sage", "mono", "navy"],
        default="sage",
        help="컬러 팔레트 (기본: sage)"
    )

    args = parser.parse_args()

    # JSON 설정 파일 로드
    config_path = os.path.expanduser(args.config)
    if not os.path.exists(config_path):
        print(f"Error: Config file not found: {config_path}", file=sys.stderr)
        sys.exit(1)

    with open(config_path, 'r', encoding='utf-8') as f:
        config = json.load(f)

    # 팔레트 오버라이드 (config에서 지정된 경우)
    palette_name = config.get("palette", args.palette)

    # PPTX 생성
    try:
        output_path = generate_pptx(config, args.output, palette_name)
        print(f"PPTX 생성 완료: {output_path}")
        print(f"슬라이드 수: {len(config.get('slides', []))}")
        print(f"팔레트: {palette_name}")
    except Exception as e:
        print(f"Error: PPTX 생성 실패: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
